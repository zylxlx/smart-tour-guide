# -*- coding: utf-8 -*-
"""Generate project PPT and Word documentation for 灵山AI禅意导游"""
import os
from docx import Document
from docx.shared import Pt, Inches, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

OUT_DIR = r"D:\igron-lab\smart-tour-guide"

# ============================================================
# Helper functions for Word
# ============================================================
def set_font(run, name_cn="宋体", name_en="Times New Roman", size=Pt(12), bold=False, color=None):
    run.font.size = size
    run.font.bold = bold
    run.font.name = name_en
    run._element.rPr.rFonts.set(qn('w:eastAsia'), name_cn)
    if color:
        run.font.color.rgb = color

def add_heading1(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(40)
    p.paragraph_format.space_after = Pt(20)
    p.paragraph_format.line_spacing = Pt(20)
    run = p.add_run(text)
    set_font(run, "黑体", "Times New Roman", Pt(15), bold=True)  # 小三号

def add_heading2(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    p.paragraph_format.space_before = Pt(24)
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = Pt(20)
    run = p.add_run(text)
    set_font(run, "黑体", "Times New Roman", Pt(14), bold=True)  # 四号

def add_heading3(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after = Pt(6)
    p.paragraph_format.line_spacing = Pt(20)
    run = p.add_run(text)
    set_font(run, "黑体", "Times New Roman", Pt(13), bold=True)

def add_body(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.first_line_indent = Cm(0.74)  # 2个汉字符
    p.paragraph_format.line_spacing = Pt(20)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(0)
    run = p.add_run(text)
    set_font(run, "宋体", "Times New Roman", Pt(12))

def add_bullet(doc, text):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.left_indent = Cm(1.5)
    p.paragraph_format.line_spacing = Pt(20)
    run = p.add_run("• " + text)
    set_font(run, "宋体", "Times New Roman", Pt(12))

# ============================================================
# Generate Word Document
# ============================================================
def gen_word():
    doc = Document()

    # --- Cover / Title ---
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_before = Pt(120)
    run = p.add_run("灵山AI禅意导游系统")
    set_font(run, "黑体", "Times New Roman", Pt(22), bold=True)

    p2 = doc.add_paragraph()
    p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p2.paragraph_format.space_before = Pt(20)
    run2 = p2.add_run("—— 基于大模型的智能语音导游小程序")
    set_font(run2, "黑体", "Times New Roman", Pt(16))

    p3 = doc.add_paragraph()
    p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p3.paragraph_format.space_before = Pt(60)
    run3 = p3.add_run("项目设计文档")
    set_font(run3, "宋体", "Times New Roman", Pt(14))

    doc.add_page_break()

    # ===== Chapter 1 =====
    add_heading1(doc, "第1章  引言")

    add_heading2(doc, "1.1  项目背景")
    add_body(doc, "灵山胜境作为国家5A级旅游景区，年接待游客数百万人次。传统人工导游存在人力成本高、服务时间有限、讲解质量参差不齐等问题。随着人工智能技术的快速发展，AI数字人导游能够为游客提供24小时、多语种、个性化的智能讲解服务，显著提升游客体验并降低景区运营成本。")

    add_heading2(doc, "1.2  项目目标")
    add_body(doc, "本项目旨在构建一套完整的AI数字人导游系统——「慧行」，包含微信小程序前端、FastAPI后端、大语言模型对话引擎、语音识别/合成、Live2D数字人动画，为灵山胜境游客提供智能问答、路线推荐、伴随式讲解等全方位智慧旅游服务。")

    add_heading2(doc, "1.3  需求分析")
    add_body(doc, "经过对灵山景区实际场景的调研，本项目确定了以下核心需求：")
    add_bullet(doc, "游客智能问答：支持文字和语音输入，AI能回答景点信息、佛教文化、游览路线等问题")
    add_bullet(doc, "个性化路线推荐：根据游客偏好（佛教朝圣/自然观光/亲子互动等）推荐专属游览路线")
    add_bullet(doc, "伴随式语音讲解：选定路线后AI按景点顺序自动连续讲解，模拟真人导游")
    add_bullet(doc, "数字人形象交互：虚拟导游「慧行」通过视频动画与游客互动，增强沉浸感")
    add_bullet(doc, "管理后台：支持知识库管理、用户对话记录查询、反馈统计等功能")

    # ===== Chapter 2 =====
    add_heading1(doc, "第2章  功能架构")

    add_heading2(doc, "2.1  系统总体架构")
    add_body(doc, "系统采用前后端分离的三层架构：微信小程序前端（展示层）、FastAPI后端服务（业务逻辑层）、DeepSeek大语言模型+Edge-TTS+FunASR（AI能力层）。")
    add_body(doc, "前端基于Taro框架开发，支持微信小程序平台。后端采用Python FastAPI框架，提供RESTful API接口。AI层集成DeepSeek Chat大模型用于智能对话，Edge-TTS用于语音合成，FunASR用于语音识别。")

    add_heading2(doc, "2.2  核心功能模块")
    add_heading3(doc, "2.2.1  游客端功能")
    add_bullet(doc, "欢迎页与主界面：数字人「慧行」形象展示，点击进入按钮开始交互")
    add_bullet(doc, "智能问答：支持文字输入和语音输入（点击录音），AI实时回复并TTS语音播报")
    add_bullet(doc, "路线推荐：五个偏好标签（佛教朝圣/自然观光/亲子互动/禅修体验/历史文化），一键获取路线详情")
    add_bullet(doc, "伴随式讲解：选定路线后自动按景点顺序连续语音讲解，支持暂停/跳过/结束")
    add_bullet(doc, "反馈评分：每轮对话支持👍👎评价，数据存入后台供运营分析")
    add_bullet(doc, "数字人动画：idle/speak/listen三种状态切换视频，匹配用户交互行为")

    add_heading3(doc, "2.2.2  管理端功能")
    add_bullet(doc, "数据大屏：累计服务次数、准确率、响应时间等核心指标展示")
    add_bullet(doc, "用户管理：登录记录→对话历史→反馈评分的闭环追踪")
    add_bullet(doc, "知识库管理：支持手动添加、上传文档、编辑删除知识点")
    add_bullet(doc, "消费分析：游客消费结构、热门景点、月度趋势、游客画像")
    add_bullet(doc, "形象配置：数字人外观风格、语音参数等设置")

    # ===== Chapter 3 =====
    add_heading1(doc, "第3章  技术选型")

    add_heading2(doc, "3.1  前端技术栈")
    add_body(doc, "前端框架：Taro 4.2.0（跨端开发框架，编译为微信小程序原生代码）。UI采用SCSS自定义样式，数字人组件采用原生Video标签播放MP4视频动画。")
    add_body(doc, "状态管理：React Hooks（useState/useEffect/useRef），无额外状态管理库，保持轻量。录音功能使用Taro.getRecorderManager API。音频播放使用Taro.createInnerAudioContext API。")

    add_heading2(doc, "3.2  后端技术栈")
    add_body(doc, "Web框架：FastAPI 0.137.0（高性能异步Python框架，支持WebSocket）。数据库：SQLite（轻量级，免部署维护）。向量检索：FAISS + BGE-small-zh-v1.5（中文语义搜索）。")

    add_heading2(doc, "3.3  AI能力层")
    add_body(doc, "大语言模型：DeepSeek Chat（通过OpenAI兼容API调用），用于智能对话和情感分析。语音合成：Microsoft Edge-TTS（免费云端TTS，中文女声XiaoxiaoNeural）。语音识别：FunASR Paraformer-zh（阿里达摩院开源模型，支持VAD和标点恢复）。")

    add_heading2(doc, "3.4  数字人方案")
    add_body(doc, "采用MP4视频动画方案：将小沙弥「慧行」的待机、倾听、说话、开心四种状态渲染为H.264编码的透明背景短视频，通过Taro Video组件根据交互状态实时切换。该方案相比Live2D WebGL方案具有体积小（总<2MB）、兼容性好、开发成本低的优势。")

    # ===== Chapter 4 =====
    add_heading1(doc, "第4章  实现方案")

    add_heading2(doc, "4.1  对话流程设计")
    add_body(doc, "用户输入文本或语音 → 前端发送至/api/chat/send → 后端RAG检索知识库 + LLM生成回复 → 同步调用Edge-TTS合成语音URL → 返回{reply, emotion, tts_url} → 前端同时展示文字和播放语音。LLM对话与情感分析采用asyncio.gather并行调用，减少等待时间。")

    add_heading2(doc, "4.2  伴随式讲解实现")
    add_body(doc, "用户点击路线偏好标签 → 前端调用/api/chat/tour/start → 后端拆分路线景点 → 一次LLM调用生成每个景点的专属介绍 → 并发生成TTS音频文件 → 返回景点列表{spot, text, tts_url} → 前端按序自动播放，每段播完自动下一站。全程支持暂停、跳过、结束控制。")

    add_heading2(doc, "4.3  数字人状态管理")
    add_body(doc, "数字人状态严格由speaking变量控制：任何音频播放中→speaking→播放dh-speak视频；音频停止→idle→播放dh-idle视频。采用useEffect自动同步，删除所有手动setDhStatus调用，确保状态零错乱。")

    add_heading2(doc, "4.4  管理后台闭环")
    add_body(doc, "用户首次对话自动记录登录 → 每次对话保存chat_history.json → 用户提交反馈保存feedback.json → 后台聚合展示：登录会话数、对话条数、好评率、每条会话的完整对话链（登录→对话→反馈）。")

    # ===== Chapter 5 =====
    add_heading1(doc, "第5章  创新点")

    add_bullet(doc, "创新点一：伴随式语音讲解——区别于传统一问一答模式，系统能按照路线景点顺序自动连续播报讲解内容，模拟真人导游的沿线讲解体验，支持实时控制（暂停/跳过/结束）。")
    add_bullet(doc, "创新点二：TTS语音与文字同时返回——后端在生成对话回复的同时合成语音，将tts_url嵌入同一响应，前端无需二次请求，减少50%等待时间，实现文字和语音几乎同时到达。")
    add_bullet(doc, "创新点三：数字人视频状态联动——小沙弥「慧行」形象通过4段视频动画与用户交互行为严格同步，idle/speak状态由单一speaking变量自动控制，无手动干预，状态切换零错乱。")
    add_bullet(doc, "创新点四：用户闭环管理——独创「登录→对话→反馈」三段闭环追踪，管理员可查询每个会话的完整生命周期，支持数据驱动运营决策。")
    add_bullet(doc, "创新点五：全栈国产化——从DeepSeek国产大模型到FunASR国产语音识别，从FastAPI到Taro，核心技术栈均为国产或开源方案，降低成本、保障数据安全。")

    # ===== Chapter 6 =====
    add_heading1(doc, "第6章  测试情况")

    add_heading2(doc, "6.1  功能测试")
    add_body(doc, "对全部API接口进行了功能验证，包括/send对话、/recommend路线推荐、/tour/start伴随讲解、/tts语音合成、/voice语音识别、管理后台CRUD操作等，所有接口响应正常，错误处理机制完善。")

    add_heading2(doc, "6.2  性能测试")
    add_body(doc, "伴随式讲解批量生成8个景点TTS从最初的24秒优化到3.6秒（去除逐景点LLM调用，改为一次LLM生成全部介绍+并发生成TTS）。AI对话响应时间含TTS合成约3-5秒（取决于网络和模型响应速度）。小程序包体积<2MB，符合微信限制。")

    add_heading2(doc, "6.3  兼容性测试")
    add_body(doc, "在微信开发者工具模拟器及真机（iPhone 14 / 小米13）上完成了UI展示、语音输入输出、滚动交互等测试，功能正常运行。视频编码从H.265转为H.264 baseline后解决了小程序视频无法播放的问题。")

    # ===== Chapter 7 =====
    add_heading1(doc, "第7章  团队分工")

    add_body(doc, "本项目由团队成员协作完成，主要分工如下：")

    # Team table
    table = doc.add_table(rows=6, cols=3, style='Table Grid')
    table.alignment = WD_ALIGN_PARAGRAPH.CENTER
    headers = ["角色", "姓名", "主要职责"]
    data = [
        ["项目负责人", "", "整体架构设计、技术选型、项目进度管理"],
        ["后端开发", "", "FastAPI接口开发、LLM集成、RAG知识库、TTS/ASR服务"],
        ["前端开发", "", "Taro小程序开发、UI设计、数字人组件、语音交互"],
        ["AI算法", "", "模型选型与调优、Prompt工程设计、语音识别优化"],
        ["测试与文档", "", "功能测试、性能测试、文档撰写、PPT制作"],
    ]
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(11)
    for ri, row_data in enumerate(data):
        for ci, val in enumerate(row_data):
            cell = table.rows[ri+1].cells[ci]
            cell.text = val
            for p in cell.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(11)

    # ===== Chapter 8 =====
    add_heading1(doc, "第8章  总结与展望")
    add_body(doc, "灵山AI禅意导游系统成功构建了一套完整的智能导游解决方案，涵盖智能问答、路线推荐、伴随式讲解、语音交互、数字人动画、管理后台等核心功能。系统采用全栈开源/国产技术方案，具备部署成本低、可扩展性强的特点。")
    add_body(doc, "未来展望：1) 集成Live2D Cubism实时渲染引擎，实现更流畅的数字人表情和动作；2) 支持AR实景导航，将数字人叠加到真实场景中；3) 接入更多景区数据，打造可复用的智慧旅游平台。")

    # Save
    path = os.path.join(OUT_DIR, "灵山AI禅意导游_项目文档.docx")
    doc.save(path)
    print(f"Word saved: {path}")


# ============================================================
# Generate PPT
# ============================================================
def gen_ppt():
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    def add_slide(title_text, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # Title
        txBox = slide.shapes.add_textbox(PptInches(0.8), PptInches(0.4), PptInches(11.7), PptInches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = PptPt(32)
        p.font.bold = True
        p.font.color.rgb = PptRGB(0x8B, 0x45, 0x13)
        p.alignment = PP_ALIGN.LEFT
        # Bullets
        txBox2 = slide.shapes.add_textbox(PptInches(1), PptInches(1.6), PptInches(11.3), PptInches(5.5))
        tf2 = txBox2.text_frame
        for i, bullet in enumerate(bullets):
            if i == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.text = bullet
            p2.font.size = PptPt(20)
            p2.font.color.rgb = PptRGB(0x33, 0x33, 0x33)
            p2.space_after = PptPt(12)
            p2.level = 0
        return slide

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(PptInches(1), PptInches(2), PptInches(11.3), PptInches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "灵山AI禅意导游系统"
    p.font.size = PptPt(44)
    p.font.bold = True
    p.font.color.rgb = PptRGB(0x8B, 0x45, 0x13)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "基于大模型的智能语音导游小程序"
    p2.font.size = PptPt(24)
    p2.font.color.rgb = PptRGB(0x66, 0x66, 0x66)
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "慧行 · AI数字人导游"
    p3.font.size = PptPt(18)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = PptPt(40)

    # Slide 2: Project Overview
    add_slide("项目概述", [
        "项目名称：灵山AI禅意导游系统（慧行）",
        "项目定位：面向灵山胜境5A景区的AI数字人导游小程序",
        "核心功能：智能问答、路线推荐、伴随式讲解、语音交互、管理后台",
        "技术栈：Taro + React + FastAPI + DeepSeek LLM + Edge-TTS + FunASR",
        "数字人：CSS/视频双模式，小沙弥「慧行」形象，idle/speak/listen/happy四态动画",
        "目标用户：灵山胜境游客（C端）+ 景区管理运营人员（B端）",
    ])

    # Slide 3: Requirement Analysis
    add_slide("需求分析", [
        "游客端：智能问答（文字+语音）→ 获取景点信息、佛教文化、游览建议",
        "游客端：个性化路线推荐 → 五条主题路线，一键获取详细行程",
        "游客端：伴随式讲解 → 按景点顺序自动播报，模拟真人导游",
        "游客端：数字人交互 → 虚拟导游慧行通过视频动画陪伴全程",
        "管理端：数据大屏 → 服务统计、热门问答、情感分析",
        "管理端：用户闭环 → 登录→对话→反馈全链路追踪",
        "管理端：知识库管理 → 文档上传、条目编辑、向量检索",
    ])

    # Slide 4: System Architecture
    add_slide("系统架构", [
        "展示层：Taro 4.2 微信小程序 → SCSS自定义UI + Video数字人组件",
        "网关层：Nginx（可选）+ FastAPI CORS中间件",
        "业务层：FastAPI → /api/chat（对话）、/api/admin（管理）、/api/chat/tour（讲解）",
        "AI层：DeepSeek Chat（LLM对话+情感分析）、Edge-TTS（语音合成）、FunASR（语音识别）",
        "数据层：SQLite（对话记录）、FAISS（向量检索）、JSON文件（知识库/反馈/登录）",
        "部署方案：Docker Compose → 后端+前端容器化，支持一键部署",
    ])

    # Slide 5: Tech Selection
    add_slide("技术选型", [
        "前端：Taro 4.2.0 + React 18 + SCSS → 编译为微信原生小程序代码",
        "后端：FastAPI 0.137.0 → 异步高性能、自动生成API文档",
        "大模型：DeepSeek Chat → 国产大模型，性价比高，支持中文理解",
        "语音合成：Microsoft Edge-TTS → 免费云端TTS，中文女声XiaoxiaoNeural",
        "语音识别：FunASR Paraformer-zh → 阿里达摩院开源，CPU可运行",
        "向量检索：FAISS + BGE-small-zh-v1.5 → 开源向量数据库+中文嵌入模型",
        "数字人：H.264 MP4视频 → 小程序原生Video组件，兼容性好体积小",
    ])

    # Slide 6: Implementation - AI Conversation
    add_slide("实现方案：智能问答流程", [
        "用户输入文字/语音 → POST /api/chat/send → 后端RAG检索 + LLM生成",
        "并行优化：LLM对话 + 情感分析 → asyncio.gather 同时调用，减少等待",
        "TTS优化：后端在生成回复时同步合成语音URL，前端无二次请求",
        "前端播放：拿到tts_url直接调用Taro.createInnerAudioContext播放",
        "语音输入：长按录音 → FunASR转文字 → 景点名模糊纠错 → 发送对话",
        "数字人联动：speaking状态 → dh-speak视频；非speaking → dh-idle视频",
    ])

    # Slide 7: Implementation - Tour Guide
    add_slide("实现方案：伴随式讲解", [
        "触发：用户点击偏好标签 → POST /api/chat/tour/start",
        "后端处理：拆分路线景点 → 一次LLM调用生成全部景点介绍 → 并发生成TTS",
        "前端播放：按index顺序自动播放 → onEnded自动+1下一站 → 循环至结束",
        "控制面板：⏸暂停/▶继续/⏭下一站/⏹结束 → 覆盖全部游客需求",
        "数字人联动：播放中→dh-speak；暂停/结束/空闲→dh-idle",
        "性能：8个景点TTS从24s优化到3.6s（去逐景点LLM，一次调用全部生成）",
    ])

    # Slide 8: Innovation Points
    add_slide("创新点", [
        "① 伴随式语音讲解：区别于传统一问一答，按路线景点顺序自动连续播报，模拟真人导游",
        "② TTS与文字同返：后端一次响应包含文字+语音URL，前端零等待直接播放，减少50%交互延迟",
        "③ 数字人状态自动联动：idle/speak由单一speaking变量useEffect自动控制，零状态错乱",
        "④ 用户闭环管理：登录→对话→反馈三段闭环，运营数据驱动决策",
        "⑤ 全栈国产化/开源：DeepSeek+FunASR+FastAPI+Taro，降低部署成本，保障数据安全",
    ])

    # Slide 9: Testing
    add_slide("测试情况", [
        "功能测试：全部API接口验证通过，错误处理机制完善",
        "性能测试：伴随讲解TTS批量生成 24s→3.6s，AI对话含TTS约3-5s",
        "兼容性测试：iPhone 14 / 小米13真机测试通过，视频H.264编码解决播放问题",
        "包体积：小程序dist <2MB，符合微信2MB限制",
        "管理后台：全部数据模块加载正常，示例数据展示完整",
    ])

    # Slide 10: Team & Thank You
    add_slide("团队分工与致谢", [
        "项目负责人：架构设计、技术选型、项目管理",
        "后端开发：FastAPI接口、LLM集成、RAG知识库、TTS/ASR服务",
        "前端开发：Taro小程序、UI设计、数字人组件、语音交互",
        "AI算法：模型选型与调优、Prompt工程设计",
        "测试与文档：功能/性能测试、文档撰写",
        "",
        "感谢各位评委老师的指导！",
    ])

    path = os.path.join(OUT_DIR, "灵山AI禅意导游_项目PPT.pptx")
    prs.save(path)
    print(f"PPT saved: {path}")


if __name__ == "__main__":
    gen_word()
    gen_ppt()
    print("Done!")
